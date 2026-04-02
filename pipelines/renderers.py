from __future__ import annotations

import logging
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor as DocxRGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptxInches
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

from pipelines.schemas import DeckArtifact, FieldGuideArtifact, HandoutArtifact, LMSArtifact, ManualArtifact


LOGGER = logging.getLogger(__name__)

ARIZONA_RED = RGBColor(171, 24, 39)
ARIZONA_GOLD = RGBColor(245, 187, 66)
ARIZONA_NAVY = RGBColor(22, 41, 75)
ARIZONA_SAND = RGBColor(244, 236, 219)
PALE_CARD = RGBColor(255, 249, 240)
SLIDE_KICKERS = {
    "Why We're Here": "Purpose",
    "What Is the Rule of Law?": "Definition",
    "Why It Matters": "Why It Matters",
    "What Happens When It Weakens": "When Process Fails",
    "Why Lawyers Must Lead": "Professional Duty",
    "The Ambassador Role": "Ambassador Role",
    "Communicating Effectively": "Communication",
    "What Citizens Can Do": "Public Role",
    "Arizona-Specific Resources": "Arizona Resources",
    "Call to Action": "Next Step",
    "Closing / Additional Resources": "Continue the Work",
}


def _set_document_base_styles(document: Document) -> None:
    section = document.sections[0]
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.8)
    section.right_margin = Inches(0.8)

    normal = document.styles["Normal"]
    normal.font.name = "Aptos"
    normal.font.size = Pt(10.5)
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "Aptos")

    heading1 = document.styles["Heading 1"]
    heading1.font.name = "Aptos Display"
    heading1.font.size = Pt(16)
    heading1.font.bold = True
    heading1.font.color.rgb = DocxRGBColor(22, 41, 75)
    heading1._element.rPr.rFonts.set(qn("w:eastAsia"), "Aptos Display")

    heading2 = document.styles["Heading 2"]
    heading2.font.name = "Aptos"
    heading2.font.size = Pt(12)
    heading2.font.bold = True
    heading2.font.color.rgb = DocxRGBColor(171, 24, 39)
    heading2._element.rPr.rFonts.set(qn("w:eastAsia"), "Aptos")


def _add_rule(paragraph) -> None:
    p_pr = paragraph._p.get_or_add_pPr()
    p_bdr = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), "8")
    bottom.set(qn("w:space"), "1")
    bottom.set(qn("w:color"), "AB1827")
    p_bdr.append(bottom)
    p_pr.append(p_bdr)


def write_manual_docx(manual: ManualArtifact, path: Path) -> None:
    document = Document()
    _set_document_base_styles(document)
    document.core_properties.title = manual.title
    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run(manual.title)
    run.bold = True
    run.font.size = Pt(24)
    run.font.name = "Aptos Display"
    subtitle = document.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle_run = subtitle.add_run("Arizona Rule of Law Ambassador Program")
    subtitle_run.font.size = Pt(11)
    subtitle_run.font.color.rgb = DocxRGBColor(171, 24, 39)
    divider = document.add_paragraph()
    _add_rule(divider)
    for section in manual.sections:
        document.add_heading(section.title, level=1)
        for block in section.body_markdown.split("\n\n"):
            document.add_paragraph(block.replace("**", ""))
        for callout in section.callouts:
            document.add_paragraph(callout, style="Intense Quote")
        for checklist in section.checklists:
            document.add_paragraph(checklist.title, style="Heading 2")
            for item in checklist.items:
                document.add_paragraph(item, style="List Bullet")
        for quote in section.pull_quotes:
            document.add_paragraph(quote, style="Quote")
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def write_field_guide_docx(field_guide: FieldGuideArtifact, path: Path) -> None:
    document = Document()
    _set_document_base_styles(document)
    document.add_heading("Arizona Rule of Law Ambassador Speaker Field Guide", level=0)
    for version in field_guide.versions:
        document.add_heading(f"{version.minutes}-Minute Version", level=1)
        for block in version.blocks:
            document.add_heading(f"{block.title} ({block.minutes} min)", level=2)
            document.add_paragraph(block.script)
            document.add_paragraph(f"Transition: {block.transition}")
            document.add_paragraph(f"Audience engagement: {block.engagement_prompt}")
            document.add_paragraph(f"Speaker tip: {block.speaker_tip}")
            if block.likely_questions:
                document.add_paragraph("Likely questions", style="Heading 3")
                for question in block.likely_questions:
                    document.add_paragraph(question, style="List Bullet")
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def write_handout_docx(handout: HandoutArtifact, path: Path) -> None:
    document = Document()
    _set_document_base_styles(document)
    document.add_heading(handout.title, level=0)
    for section in handout.sections:
        document.add_heading(section.title, level=1)
        document.add_paragraph(section.body)
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def write_lms_docx(lms: LMSArtifact, path: Path) -> None:
    document = Document()
    _set_document_base_styles(document)
    document.add_heading("Arizona Rule of Law Ambassador LMS Modules", level=0)
    for module in lms.modules:
        document.add_heading(module.title, level=1)
        document.add_paragraph(module.overview)
        document.add_paragraph("Learning objectives", style="Heading 2")
        for objective in module.learning_objectives:
            document.add_paragraph(objective, style="List Bullet")
        document.add_paragraph(module.lesson_body_markdown.replace("**", ""))
        document.add_paragraph("Key takeaways", style="Heading 2")
        for takeaway in module.key_takeaways:
            document.add_paragraph(takeaway, style="List Bullet")
        document.add_paragraph("Reflection questions", style="Heading 2")
        for question in module.reflection_questions:
            document.add_paragraph(question, style="List Bullet")
        document.add_paragraph("Knowledge check", style="Heading 2")
        for quiz in module.quiz_questions:
            document.add_paragraph(quiz.question, style="List Number")
            for choice in quiz.choices:
                document.add_paragraph(choice, style="List Bullet 2")
            document.add_paragraph(f"Answer: {quiz.answer}")
            document.add_paragraph(f"Rationale: {quiz.rationale}")
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def write_simple_pdf(title: str, sections: list[tuple[str, str]], path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="SectionHeader", parent=styles["Heading2"], textColor=colors.HexColor("#16294B"), spaceAfter=8))
    body_style = ParagraphStyle("BodyCopy", parent=styles["BodyText"], fontSize=10, leading=14, spaceAfter=8)
    story = [Paragraph(title, styles["Title"]), Spacer(1, 12)]
    for section_title, body in sections:
        story.append(Paragraph(section_title, styles["SectionHeader"]))
        for block in body.split("\n\n"):
            story.append(Paragraph(block.replace("\n", "<br/>"), body_style))
        story.append(Spacer(1, 8))
    document = SimpleDocTemplate(str(path), pagesize=letter, rightMargin=54, leftMargin=54, topMargin=54, bottomMargin=54)
    document.build(story)


def write_deck_pptx(deck: DeckArtifact, path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = PptxInches(13.333)
    presentation.slide_height = PptxInches(7.5)
    blank_layout = presentation.slide_layouts[6]
    for slide_artifact in deck.slides:
        slide = presentation.slides.add_slide(blank_layout)
        _paint_slide_background(slide, slide_artifact.slide_number)
        if slide_artifact.slide_number == 1:
            _render_title_slide(slide, slide_artifact)
        elif slide_artifact.title.startswith("Core Principle"):
            _render_principle_slide(slide, slide_artifact)
        elif slide_artifact.title == "What Happens When It Weakens":
            _render_warning_slide(slide, slide_artifact)
        elif "Resources" in slide_artifact.title:
            _render_resources_slide(slide, slide_artifact)
        elif slide_artifact.title == "Call to Action":
            _render_action_slide(slide, slide_artifact)
        else:
            _render_standard_slide(slide, slide_artifact)
        try:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_artifact.speaker_notes
        except Exception:  # noqa: BLE001
            LOGGER.info("Speaker notes could not be embedded for slide %s.", slide_artifact.slide_number)
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def _paint_slide_background(slide, slide_number: int) -> None:
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = ARIZONA_SAND
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, PptxInches(13.333), PptxInches(0.35))
    band.fill.solid()
    band.fill.fore_color.rgb = ARIZONA_RED
    band.line.fill.background()
    sun = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, PptxInches(10.4), PptxInches(0.45), PptxInches(1.7), PptxInches(1.7))
    sun.fill.solid()
    sun.fill.fore_color.rgb = ARIZONA_GOLD
    sun.fill.transparency = 0.28
    sun.line.fill.background()
    horizon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, PptxInches(8.3), PptxInches(5.95), PptxInches(5.5), PptxInches(1.3))
    horizon.fill.solid()
    horizon.fill.fore_color.rgb = ARIZONA_RED
    horizon.fill.transparency = 0.18
    horizon.line.fill.background()
    side = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PptxInches(12.78), 0, PptxInches(0.55), PptxInches(7.5))
    side.fill.solid()
    side.fill.fore_color.rgb = ARIZONA_NAVY
    side.fill.transparency = 0.05
    side.line.fill.background()
    marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, PptxInches(12.9), PptxInches(0.55), PptxInches(0.28), PptxInches(0.28))
    marker.fill.solid()
    marker.fill.fore_color.rgb = ARIZONA_GOLD
    marker.line.fill.background()
    number_box = slide.shapes.add_textbox(PptxInches(12.62), PptxInches(6.7), PptxInches(0.55), PptxInches(0.35))
    number_p = number_box.text_frame.paragraphs[0]
    number_p.text = f"{slide_number:02d}"
    number_p.font.name = "Aptos Display"
    number_p.font.size = Pt(10)
    number_p.font.color.rgb = ARIZONA_SAND
    number_p.alignment = PP_ALIGN.CENTER


def _render_title_slide(slide, slide_artifact) -> None:
    kicker = slide.shapes.add_textbox(PptxInches(0.9), PptxInches(0.9), PptxInches(4.0), PptxInches(0.35))
    kicker_p = kicker.text_frame.paragraphs[0]
    kicker_p.text = "Arizona Rule of Law Ambassador Program"
    kicker_p.font.name = "Aptos"
    kicker_p.font.size = Pt(12)
    kicker_p.font.color.rgb = ARIZONA_RED

    title_box = slide.shapes.add_textbox(PptxInches(0.85), PptxInches(1.35), PptxInches(7.0), PptxInches(2.2))
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = slide_artifact.title
    title_p.font.name = "Aptos Display"
    title_p.font.size = Pt(30)
    title_p.font.bold = True
    title_p.font.color.rgb = ARIZONA_NAVY

    if slide_artifact.body_lines:
        headline_p = title_box.text_frame.add_paragraph()
        headline_p.text = slide_artifact.body_lines[0]
        headline_p.font.name = "Aptos Display"
        headline_p.font.size = Pt(22)
        headline_p.font.bold = True
        headline_p.font.color.rgb = ARIZONA_RED
        headline_p.space_before = Pt(10)
    if len(slide_artifact.body_lines) > 1:
        support_p = title_box.text_frame.add_paragraph()
        support_p.text = slide_artifact.body_lines[1]
        support_p.font.name = "Aptos"
        support_p.font.size = Pt(15)
        support_p.font.color.rgb = ARIZONA_NAVY
        support_p.space_before = Pt(8)

    subtitle_p = title_box.text_frame.add_paragraph()
    subtitle_p.text = "Civic education materials for lawyers speaking with schools, community groups, and civic organizations."
    subtitle_p.font.name = "Aptos"
    subtitle_p.font.size = Pt(14)
    subtitle_p.font.color.rgb = ARIZONA_NAVY
    subtitle_p.space_before = Pt(16)

    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptxInches(8.1), PptxInches(1.2), PptxInches(3.9), PptxInches(2.35))
    panel.fill.solid()
    panel.fill.fore_color.rgb = ARIZONA_NAVY
    panel.line.color.rgb = ARIZONA_GOLD
    panel.text_frame.word_wrap = True
    panel.text_frame.paragraphs[0].text = "Program Focus"
    panel.text_frame.paragraphs[0].font.name = "Aptos Display"
    panel.text_frame.paragraphs[0].font.size = Pt(18)
    panel.text_frame.paragraphs[0].font.color.rgb = ARIZONA_SAND
    for line in [
        "Nonpartisan civic education",
        "Plain-English legal communication",
        "Arizona-rooted public service",
    ]:
        p = panel.text_frame.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(14)
        p.font.color.rgb = ARIZONA_SAND
        p.level = 0

    _add_footer(slide, slide_artifact.slide_number, emphasis="Professional civic education")


def _render_standard_slide(slide, slide_artifact) -> None:
    _add_kicker(slide, SLIDE_KICKERS.get(slide_artifact.title))
    _add_title(slide, slide_artifact.title)
    _add_left_rule(slide)
    body_card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptxInches(0.85), PptxInches(1.55), PptxInches(6.4), PptxInches(4.6))
    body_card.fill.solid()
    body_card.fill.fore_color.rgb = PALE_CARD
    body_card.line.color.rgb = ARIZONA_GOLD
    body_tf = body_card.text_frame
    body_tf.word_wrap = True
    body_tf.margin_left = PptxInches(0.28)
    body_tf.margin_right = PptxInches(0.28)
    body_tf.margin_top = PptxInches(0.22)
    for idx, line in enumerate(slide_artifact.body_lines):
        p = body_tf.paragraphs[0] if idx == 0 else body_tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos Display" if idx == 0 else "Aptos"
        p.font.size = Pt(24 if idx == 0 else 18)
        p.font.bold = idx == 0
        p.font.color.rgb = ARIZONA_NAVY
        p.space_after = Pt(10)

    _add_focus_panel(slide, slide_artifact, x=7.55, y=1.55, w=4.5, h=3.95)
    _add_footer(slide, slide_artifact.slide_number)


def _render_principle_slide(slide, slide_artifact) -> None:
    _add_kicker(slide, "Core Principle")
    _add_title(slide, slide_artifact.title)
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, PptxInches(0.95), PptxInches(1.55), PptxInches(1.1), PptxInches(1.1))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ARIZONA_RED
    badge.line.fill.background()
    badge_p = badge.text_frame.paragraphs[0]
    badge_p.text = str(slide_artifact.slide_number - 4)
    badge_p.font.name = "Aptos Display"
    badge_p.font.size = Pt(24)
    badge_p.font.bold = True
    badge_p.font.color.rgb = ARIZONA_SAND
    badge_p.alignment = PP_ALIGN.CENTER

    principle = slide.shapes.add_textbox(PptxInches(2.25), PptxInches(1.65), PptxInches(4.7), PptxInches(1.0))
    p = principle.text_frame.paragraphs[0]
    p.text = slide_artifact.body_lines[0] if slide_artifact.body_lines else slide_artifact.title
    p.font.name = "Aptos Display"
    p.font.size = Pt(27)
    p.font.bold = True
    p.font.color.rgb = ARIZONA_NAVY

    definition = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptxInches(0.95), PptxInches(2.8), PptxInches(5.9), PptxInches(2.3))
    definition.fill.solid()
    definition.fill.fore_color.rgb = PALE_CARD
    definition.line.color.rgb = ARIZONA_GOLD
    definition.text_frame.margin_left = PptxInches(0.24)
    definition.text_frame.margin_right = PptxInches(0.24)
    definition.text_frame.margin_top = PptxInches(0.18)
    d_p = definition.text_frame.paragraphs[0]
    d_p.text = slide_artifact.body_lines[1] if len(slide_artifact.body_lines) > 1 else ""
    d_p.font.name = "Aptos"
    d_p.font.size = Pt(20)
    d_p.font.color.rgb = ARIZONA_NAVY

    _add_focus_panel(slide, slide_artifact, x=7.5, y=1.55, w=4.6, h=3.8, label="Why it matters")
    _add_footer(slide, slide_artifact.slide_number, emphasis="Four core principles")


def _render_resources_slide(slide, slide_artifact) -> None:
    _add_kicker(slide, SLIDE_KICKERS.get(slide_artifact.title))
    _add_title(slide, slide_artifact.title)
    headings = slide_artifact.body_lines[:3]
    x_positions = [0.9, 4.4, 7.9]
    for x, heading in zip(x_positions, headings):
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptxInches(x), PptxInches(2.0), PptxInches(3.0), PptxInches(2.7))
        card.fill.solid()
        card.fill.fore_color.rgb = PALE_CARD
        card.line.color.rgb = ARIZONA_GOLD
        tf = card.text_frame
        tf.margin_left = PptxInches(0.18)
        tf.margin_right = PptxInches(0.18)
        tf.margin_top = PptxInches(0.18)
        title_p = tf.paragraphs[0]
        title_p.text = heading
        title_p.font.name = "Aptos Display"
        title_p.font.size = Pt(18)
        title_p.font.bold = True
        title_p.font.color.rgb = ARIZONA_NAVY
        body_p = tf.add_paragraph()
        body_p.text = "Add approved links, contact details, or local event information here."
        body_p.font.name = "Aptos"
        body_p.font.size = Pt(13)
        body_p.font.color.rgb = ARIZONA_NAVY
    if slide_artifact.optional_quote:
        _add_quote_band(slide, slide_artifact.optional_quote, y=5.25)
    _add_footer(slide, slide_artifact.slide_number, emphasis="Local partners and follow-up")


def _render_action_slide(slide, slide_artifact) -> None:
    _add_kicker(slide, SLIDE_KICKERS.get(slide_artifact.title))
    _add_title(slide, slide_artifact.title)
    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptxInches(0.9), PptxInches(1.7), PptxInches(5.7), PptxInches(3.8))
    left.fill.solid()
    left.fill.fore_color.rgb = ARIZONA_NAVY
    left.line.color.rgb = ARIZONA_GOLD
    tf = left.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(slide_artifact.body_lines[:2]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos Display" if idx == 0 else "Aptos"
        p.font.size = Pt(24 if idx == 0 else 18)
        p.font.bold = idx == 0
        p.font.color.rgb = ARIZONA_SAND
        p.space_after = Pt(12)
    for line in slide_artifact.body_lines[2:]:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(16)
        p.font.color.rgb = ARIZONA_SAND

    right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptxInches(7.0), PptxInches(1.7), PptxInches(5.0), PptxInches(3.8))
    right.fill.solid()
    right.fill.fore_color.rgb = PALE_CARD
    right.line.color.rgb = ARIZONA_RED
    rtf = right.text_frame
    rtf.margin_left = PptxInches(0.2)
    rtf.margin_right = PptxInches(0.2)
    rtf.margin_top = PptxInches(0.2)
    head = rtf.paragraphs[0]
    head.text = "Speaker Prompt"
    head.font.name = "Aptos Display"
    head.font.size = Pt(18)
    head.font.bold = True
    head.font.color.rgb = ARIZONA_NAVY
    body = rtf.add_paragraph()
    body.text = "Invite the audience to leave with one practical commitment: explain the principle, defend fair process, and point others toward trusted civic institutions."
    body.font.name = "Aptos"
    body.font.size = Pt(15)
    body.font.color.rgb = ARIZONA_NAVY
    if slide_artifact.optional_quote:
        _add_quote_band(slide, slide_artifact.optional_quote, y=5.8)
    _add_footer(slide, slide_artifact.slide_number, emphasis="Closing commitment")


def _render_warning_slide(slide, slide_artifact) -> None:
    _add_kicker(slide, SLIDE_KICKERS.get(slide_artifact.title))
    _add_title(slide, slide_artifact.title)
    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptxInches(0.85), PptxInches(1.55), PptxInches(5.9), PptxInches(4.5))
    left.fill.solid()
    left.fill.fore_color.rgb = ARIZONA_NAVY
    left.line.color.rgb = ARIZONA_RED
    tf = left.text_frame
    tf.margin_left = PptxInches(0.24)
    tf.margin_right = PptxInches(0.24)
    tf.margin_top = PptxInches(0.2)
    tf.word_wrap = True
    for idx, line in enumerate(slide_artifact.body_lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos Display" if idx == 0 else "Aptos"
        p.font.size = Pt(24 if idx == 0 else 18)
        p.font.bold = idx == 0
        p.font.color.rgb = ARIZONA_SAND
        p.space_after = Pt(12)
    _add_focus_panel(slide, slide_artifact, x=7.25, y=1.55, w=4.8, h=4.5, label="Why this matters")
    _add_footer(slide, slide_artifact.slide_number, emphasis="Protecting public trust")


def _add_title(slide, text: str) -> None:
    title_box = slide.shapes.add_textbox(PptxInches(0.85), PptxInches(0.68), PptxInches(8.6), PptxInches(0.75))
    p = title_box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos Display"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ARIZONA_NAVY


def _add_kicker(slide, text: str | None) -> None:
    if not text:
        return
    kicker = slide.shapes.add_textbox(PptxInches(0.85), PptxInches(0.42), PptxInches(3.0), PptxInches(0.25))
    p = kicker.text_frame.paragraphs[0]
    p.text = text.upper()
    p.font.name = "Aptos"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = ARIZONA_RED


def _add_left_rule(slide) -> None:
    rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PptxInches(0.65), PptxInches(1.55), PptxInches(0.14), PptxInches(4.65))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ARIZONA_RED
    rule.line.fill.background()


def _add_focus_panel(slide, slide_artifact, x: float, y: float, w: float, h: float, label: str = "Takeaway") -> None:
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = PALE_CARD
    panel.line.color.rgb = ARIZONA_GOLD
    tf = panel.text_frame
    tf.word_wrap = True
    tf.margin_left = PptxInches(0.2)
    tf.margin_right = PptxInches(0.2)
    tf.margin_top = PptxInches(0.18)
    p = tf.paragraphs[0]
    p.text = label
    p.font.name = "Aptos Display"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = ARIZONA_NAVY
    for term in _focus_terms(slide_artifact.body_lines):
        chip = tf.add_paragraph()
        chip.text = term
        chip.font.name = "Aptos"
        chip.font.size = Pt(14)
        chip.font.bold = True
        chip.font.color.rgb = ARIZONA_RED
        chip.space_before = Pt(8)
    summary = tf.add_paragraph()
    summary.text = slide_artifact.speaker_notes[:220].rstrip() + ("..." if len(slide_artifact.speaker_notes) > 220 else "")
    summary.font.name = "Aptos"
    summary.font.size = Pt(13)
    summary.font.color.rgb = ARIZONA_NAVY
    summary.space_before = Pt(12)
    if slide_artifact.optional_quote:
        quote = tf.add_paragraph()
        quote.text = f"\"{slide_artifact.optional_quote}\""
        quote.font.name = "Aptos"
        quote.font.size = Pt(11)
        quote.font.italic = True
        quote.font.color.rgb = ARIZONA_NAVY
        quote.space_before = Pt(12)


def _focus_terms(body_lines: list[str]) -> list[str]:
    terms: list[str] = []
    for line in body_lines[:3]:
        cleaned = line.replace("[link]", "").strip(". ")
        if not cleaned:
            continue
        if len(cleaned) > 34:
            cleaned = cleaned[:31].rstrip() + "..."
        if cleaned not in terms:
            terms.append(cleaned)
    return terms


def _add_quote_band(slide, quote: str, y: float) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptxInches(0.95), PptxInches(y), PptxInches(11.0), PptxInches(0.72))
    band.fill.solid()
    band.fill.fore_color.rgb = ARIZONA_NAVY
    band.fill.transparency = 0.05
    band.line.color.rgb = ARIZONA_GOLD
    tf = band.text_frame
    tf.margin_left = PptxInches(0.18)
    tf.margin_right = PptxInches(0.18)
    tf.margin_top = PptxInches(0.08)
    p = tf.paragraphs[0]
    p.text = f"\"{quote}\""
    p.font.name = "Aptos"
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = ARIZONA_SAND


def _add_footer(slide, slide_number: int, emphasis: str | None = None) -> None:
    footer = slide.shapes.add_textbox(PptxInches(0.82), PptxInches(6.72), PptxInches(11.5), PptxInches(0.34))
    p = footer.text_frame.paragraphs[0]
    text = "Arizona Rule of Law Ambassador Program"
    if emphasis:
        text += f"   |   {emphasis}"
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(10)
    p.font.color.rgb = ARIZONA_NAVY
